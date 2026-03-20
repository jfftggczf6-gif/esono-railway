"""
ESONO Document Parser — Micro-service FastAPI
Parse any document (PDF, Excel, Word, PowerPoint, images, CSV) and return structured text.
Deployed independently (Railway/Render/Hetzner), called by the Lovable frontend.
"""

import os
import io
import tempfile
import traceback
from typing import Optional
from fastapi import FastAPI, UploadFile, File, HTTPException, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import uvicorn

app = FastAPI(title="ESONO Document Parser", version="1.0.0")

# CORS — allow Lovable frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, restrict to your Lovable domain
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Auth token (simple shared secret between frontend and parser)
PARSER_API_KEY = os.getenv("PARSER_API_KEY", "esono-parser-dev-key")

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB max
MAX_CHARS_PER_FILE = 50_000  # 50K chars max per file


# ═══════════════════════════════════════════════════════════════
# PDF PARSING — pymupdf (fitz) + pdfplumber for tables
# ═══════════════════════════════════════════════════════════════
def parse_pdf(file_bytes: bytes, filename: str) -> dict:
    """Parse PDF: extract text + tables. Works on text PDFs AND scanned PDFs (with OCR)."""
    import fitz  # pymupdf
    
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    total_pages = len(doc)
    max_pages = min(total_pages, 50)
    
    full_text = ""
    tables_found = 0
    
    for page_num in range(max_pages):
        page = doc[page_num]
        
        # Extract text (works for text PDFs)
        page_text = page.get_text("text")
        
        # If no text found, try OCR via pymupdf's built-in OCR
        if len(page_text.strip()) < 20:
            try:
                # Try to get text from images (OCR)
                page_text = page.get_text("text", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            except:
                pass
        
        if page_text.strip():
            full_text += f"\n--- Page {page_num + 1}/{total_pages} ---\n{page_text.strip()}\n"
    
    doc.close()
    
    # Try pdfplumber for table extraction (better for financial statements)
    try:
        import pdfplumber
        
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages[:max_pages]):
                tables = page.extract_tables()
                for table in tables:
                    if table and len(table) > 1:
                        tables_found += 1
                        full_text += f"\n┌═══ Tableau page {i+1} ═══\n"
                        for row in table:
                            cells = [str(c).strip() if c else "—" for c in row]
                            full_text += f"│ {' | '.join(cells)}\n"
                        full_text += f"└═══ Fin tableau ═══\n"
    except Exception as e:
        # pdfplumber failed — that's OK, we have the text from pymupdf
        pass
    
    if max_pages < total_pages:
        full_text += f"\n[... {total_pages - max_pages} pages supplémentaires non lues]\n"
    
    # Classify quality
    text_len = len(full_text.strip())
    if text_len > 1000:
        quality = "high"
    elif text_len > 200:
        quality = "medium"
    elif text_len > 0:
        quality = "low"
    else:
        quality = "failed"
    
    return {
        "content": full_text[:MAX_CHARS_PER_FILE],
        "pages": total_pages,
        "tables_found": tables_found,
        "quality": quality,
        "method": "pymupdf+pdfplumber",
    }


# ═══════════════════════════════════════════════════════════════
# PDF OCR — For scanned PDFs where pymupdf finds no text
# ═══════════════════════════════════════════════════════════════
def parse_pdf_ocr(file_bytes: bytes, filename: str) -> dict:
    """OCR scanned PDF pages using Tesseract."""
    import fitz
    from PIL import Image
    import pytesseract
    
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    total_pages = len(doc)
    max_pages = min(total_pages, 20)  # OCR is slower, limit to 20 pages
    
    full_text = ""
    
    for page_num in range(max_pages):
        page = doc[page_num]
        # Render page as image (300 DPI for good OCR quality)
        pix = page.get_pixmap(dpi=300)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        # OCR with Tesseract (French + English)
        page_text = pytesseract.image_to_string(img, lang="fra+eng")
        
        if page_text.strip():
            full_text += f"\n--- Page {page_num + 1}/{total_pages} (OCR) ---\n{page_text.strip()}\n"
    
    doc.close()
    
    if max_pages < total_pages:
        full_text += f"\n[... {total_pages - max_pages} pages supplémentaires non lues (OCR)]\n"
    
    text_len = len(full_text.strip())
    quality = "high" if text_len > 500 else "medium" if text_len > 100 else "low" if text_len > 0 else "failed"
    
    return {
        "content": full_text[:MAX_CHARS_PER_FILE],
        "pages": total_pages,
        "tables_found": 0,
        "quality": quality,
        "method": "tesseract_ocr",
    }


# ═══════════════════════════════════════════════════════════════
# EXCEL PARSING — openpyxl (native, preserves everything)
# ═══════════════════════════════════════════════════════════════
def parse_excel(file_bytes: bytes, filename: str) -> dict:
    """Parse Excel: all sheets, merged cells, formulas evaluated, hidden sheets."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    
    if ext == "xls":
        import xlrd
        wb = xlrd.open_workbook(file_contents=file_bytes)
        full_text = ""
        sheet_count = wb.nsheets
        for i in range(wb.nsheets):
            ws = wb.sheet_by_index(i)
            full_text += f"\n┌═══ Feuille: \"{ws.name}\" ({ws.nrows}×{ws.ncols}) ═══\n"
            for r in range(min(ws.nrows, 500)):
                cells = []
                for c in range(ws.ncols):
                    val = ws.cell_value(r, c)
                    if val == "":
                        cells.append("—")
                    elif isinstance(val, float):
                        if abs(val) >= 1000:
                            cells.append(f"{val:,.0f}".replace(",", " "))
                        elif val % 1 != 0:
                            cells.append(f"{val:.2f}")
                        else:
                            cells.append(str(int(val)))
                    else:
                        cells.append(str(val).strip()[:200])
                if all(c == "—" for c in cells):
                    continue
                if r == 0:
                    full_text += f"│ [{'] | ['.join(cells)}]\n│ {'─' * 60}\n"
                else:
                    full_text += f"│ {' | '.join(cells)}\n"
            full_text += f"└═══ Fin \"{ws.name}\" ═══\n"
        text_len = len(full_text.strip())
        quality = "high" if text_len > 500 else "medium" if text_len > 100 else "low"
        return {"content": full_text[:MAX_CHARS_PER_FILE], "sheets": sheet_count, "quality": quality, "method": "xlrd"}
    
    import openpyxl
    
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=False)
    
    full_text = ""
    sheet_count = len(wb.sheetnames)
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        
        # Detect hidden sheets
        is_hidden = ws.sheet_state == "hidden"
        hidden_label = " [FEUILLE CACHÉE]" if is_hidden else ""
        
        rows = ws.max_row or 0
        cols = ws.max_column or 0
        
        full_text += f"\n┌═══ Feuille: \"{sheet_name}\"{hidden_label} ({rows}×{cols}) ═══\n"
        
        # Detect merged cells
        merged = list(ws.merged_cells.ranges)
        if merged:
            full_text += f"│ [{len(merged)} cellules fusionnées]\n"
        
        # Extract rows
        row_count = 0
        for row in ws.iter_rows(max_row=min(rows, 500), values_only=False):
            cells = []
            for cell in row:
                val = cell.value
                if val is None:
                    cells.append("—")
                elif isinstance(val, (int, float)):
                    # Format numbers French-style
                    if abs(val) >= 1000:
                        cells.append(f"{val:,.0f}".replace(",", " ").replace(".", ","))
                    elif isinstance(val, float):
                        cells.append(f"{val:.2f}".replace(".", ","))
                    else:
                        cells.append(str(val))
                elif hasattr(val, 'strftime'):
                    cells.append(val.strftime("%d/%m/%Y"))
                else:
                    cells.append(str(val).strip()[:200])
            
            # Skip completely empty rows
            if all(c == "—" for c in cells):
                continue
            
            # First non-empty row = headers
            if row_count == 0:
                full_text += f"│ [{'] | ['.join(cells)}]\n"
                full_text += f"│ {'─' * min(len(' | '.join(cells)), 80)}\n"
            else:
                full_text += f"│ {' | '.join(cells)}\n"
            
            row_count += 1
            if row_count > 500:
                full_text += f"│ [... lignes supplémentaires non affichées]\n"
                break
        
        full_text += f"└═══ Fin \"{sheet_name}\" ═══\n"
    
    wb.close()
    
    text_len = len(full_text.strip())
    quality = "high" if text_len > 500 else "medium" if text_len > 100 else "low"
    
    return {
        "content": full_text[:MAX_CHARS_PER_FILE],
        "sheets": sheet_count,
        "quality": quality,
        "method": "openpyxl",
    }


# ═══════════════════════════════════════════════════════════════
# WORD PARSING — python-docx (preserves tables and structure)
# ═══════════════════════════════════════════════════════════════
def parse_docx(file_bytes: bytes, filename: str) -> dict:
    """Parse Word: text + tables + headers preserved."""
    from docx import Document
    
    doc = Document(io.BytesIO(file_bytes))
    
    full_text = ""
    tables_found = 0
    
    # Extract paragraphs with style detection
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        
        style = para.style.name.lower() if para.style else ""
        
        if "heading 1" in style or "titre 1" in style:
            full_text += f"\n═══ {text} ═══\n"
        elif "heading 2" in style or "titre 2" in style:
            full_text += f"\n── {text} ──\n"
        elif "heading 3" in style or "titre 3" in style:
            full_text += f"\n─ {text} ─\n"
        elif "list" in style or "bullet" in style:
            full_text += f"  • {text}\n"
        else:
            full_text += f"{text}\n"
    
    # Extract tables
    for table in doc.tables:
        tables_found += 1
        full_text += f"\n┌═══ Tableau ═══\n"
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip()[:200] if cell.text.strip() else "—" for cell in row.cells]
            if i == 0:
                full_text += f"│ [{'] | ['.join(cells)}]\n"
                full_text += f"│ {'─' * min(len(' | '.join(cells)), 80)}\n"
            else:
                full_text += f"│ {' | '.join(cells)}\n"
        full_text += f"└═══ Fin tableau ═══\n"
    
    # Extract headers and footers
    for section in doc.sections:
        header_text = ""
        for para in section.header.paragraphs:
            if para.text.strip():
                header_text += para.text.strip() + " "
        if header_text.strip():
            full_text = f"[En-tête: {header_text.strip()}]\n" + full_text
    
    text_len = len(full_text.strip())
    quality = "high" if text_len > 500 else "medium" if text_len > 100 else "low"
    
    return {
        "content": full_text[:MAX_CHARS_PER_FILE],
        "tables_found": tables_found,
        "quality": quality,
        "method": "python-docx",
    }


# ═══════════════════════════════════════════════════════════════
# POWERPOINT PARSING — python-pptx
# ═══════════════════════════════════════════════════════════════
def parse_pptx(file_bytes: bytes, filename: str) -> dict:
    """Parse PowerPoint: text from each slide + tables."""
    from pptx import Presentation
    
    prs = Presentation(io.BytesIO(file_bytes))
    
    full_text = ""
    slide_count = len(prs.slides)
    tables_found = 0
    
    for i, slide in enumerate(prs.slides):
        full_text += f"\n--- Slide {i+1}/{slide_count} ---\n"
        
        for shape in slide.shapes:
            # Text frames (titles, content)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Detect if it's a title
                        if shape.shape_type and "TITLE" in str(shape.shape_type):
                            full_text += f"## {text}\n"
                        else:
                            full_text += f"{text}\n"
            
            # Tables
            if shape.has_table:
                tables_found += 1
                table = shape.table
                full_text += f"\n┌═══ Tableau slide {i+1} ═══\n"
                for r_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip()[:100] if cell.text.strip() else "—" for cell in row.cells]
                    if r_idx == 0:
                        full_text += f"│ [{'] | ['.join(cells)}]\n"
                        full_text += f"│ {'─' * 60}\n"
                    else:
                        full_text += f"│ {' | '.join(cells)}\n"
                full_text += f"└═══ Fin tableau ═══\n"
        
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                full_text += f"[Notes: {notes[:500]}]\n"
    
    text_len = len(full_text.strip())
    quality = "high" if text_len > 500 else "medium" if text_len > 100 else "low"
    
    return {
        "content": full_text[:MAX_CHARS_PER_FILE],
        "slides": slide_count,
        "tables_found": tables_found,
        "quality": quality,
        "method": "python-pptx",
    }


# ═══════════════════════════════════════════════════════════════
# IMAGE OCR — Tesseract (free, no API needed)
# ═══════════════════════════════════════════════════════════════
def parse_image(file_bytes: bytes, filename: str) -> dict:
    """OCR an image using Tesseract."""
    from PIL import Image
    import pytesseract
    
    img = Image.open(io.BytesIO(file_bytes))
    
    # Convert to RGB if needed
    if img.mode != "RGB":
        img = img.convert("RGB")
    
    # OCR with French + English
    text = pytesseract.image_to_string(img, lang="fra+eng")
    
    text_len = len(text.strip())
    quality = "high" if text_len > 300 else "medium" if text_len > 50 else "low" if text_len > 0 else "failed"
    
    return {
        "content": text.strip()[:MAX_CHARS_PER_FILE],
        "quality": quality,
        "method": "tesseract_ocr",
    }


# ═══════════════════════════════════════════════════════════════
# CSV/TXT PARSING — with encoding detection
# ═══════════════════════════════════════════════════════════════
def parse_text_file(file_bytes: bytes, filename: str, ext: str) -> dict:
    """Parse CSV/TXT with automatic encoding detection."""
    import chardet
    
    # Detect encoding
    detected = chardet.detect(file_bytes[:10000])
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    confidence = detected.get("confidence", 0)
    
    try:
        text = file_bytes.decode(encoding)
    except:
        try:
            text = file_bytes.decode("utf-8", errors="replace")
        except:
            text = file_bytes.decode("latin-1")
    
    if ext in ("csv", "tsv"):
        # Detect separator
        first_lines = text.split("\n")[:5]
        sample = "\n".join(first_lines)
        semicolons = sample.count(";")
        commas = sample.count(",")
        tabs = sample.count("\t")
        
        if semicolons > commas and semicolons > tabs:
            sep = ";"
            sep_name = "point-virgule"
        elif tabs > commas and tabs > semicolons:
            sep = "\t"
            sep_name = "tabulation"
        else:
            sep = ","
            sep_name = "virgule"
        
        # Reformat as readable table
        lines = text.split("\n")
        formatted = f"[Encodage: {encoding} ({confidence:.0%}), Séparateur: {sep_name}]\n"
        for i, line in enumerate(lines[:500]):
            if not line.strip():
                continue
            cells = [c.strip().strip('"') for c in line.split(sep)]
            if i == 0:
                formatted += f"[{'] | ['.join(cells)}]\n{'─' * 60}\n"
            else:
                formatted += f"{' | '.join(cells)}\n"
        
        if len(lines) > 500:
            formatted += f"\n[... {len(lines) - 500} lignes supplémentaires]\n"
        
        text = formatted
    
    return {
        "content": text[:MAX_CHARS_PER_FILE],
        "encoding": encoding,
        "quality": "high" if len(text.strip()) > 100 else "low",
        "method": f"text_{ext}",
    }


# ═══════════════════════════════════════════════════════════════
# DOCUMENT CLASSIFICATION
# ═══════════════════════════════════════════════════════════════
def classify_document(filename: str, content: str) -> str:
    """Classify document type based on filename and content."""
    import re
    
    name = filename.lower()
    text = content.lower()[:3000]
    
    # By filename
    if re.search(r"bilan|compte.?r[eé]sultat|[eé]tats?.financ|syscohada|annexe.?compt", name):
        return "etats_financiers"
    if re.search(r"relev[eé]|bancaire|banque|bni|sgbci|ecobank|bceao", name):
        return "releve_bancaire"
    if re.search(r"facture|invoice|devis|bon.?commande", name):
        return "facture"
    if re.search(r"budget|pr[eé]vision|tr[eé]sorerie|forecast|plan.?financ", name):
        return "budget_previsionnel"
    if re.search(r"business.?plan|bp|plan.?affaire|[eé]tude.?march", name):
        return "business_plan"
    if re.search(r"statut|rccm|registre|commerce|attestation|fiscal|immatricul", name):
        return "document_legal"
    if re.search(r"rapport|activit[eé]|annual|annuel", name):
        return "rapport_activite"
    if re.search(r"organi|rh|personnel|emploi|poste", name):
        return "organigramme_rh"
    if re.search(r"\.(jpg|jpeg|png|gif|webp|bmp|tiff)$", name):
        return "photo_installation"
    
    # By content
    if re.search(r"total.?actif|total.?passif|capitaux.?propres|immobilisation", text):
        return "etats_financiers"
    if re.search(r"solde|d[eé]bit|cr[eé]dit|virement|pr[eé]l[eè]vement", text):
        return "releve_bancaire"
    if re.search(r"montant.?ht|montant.?ttc|tva|facture.?n", text):
        return "facture"
    if re.search(r"pr[eé]visionnel|projection|hypoth[eè]se|sc[eé]nario", text):
        return "budget_previsionnel"
    if re.search(r"proposition.?valeur|segment|canvas|march[eé].?cible", text):
        return "business_plan"
    if re.search(r"article|associ[eé]|g[eé]rant|capital.?social", text):
        return "document_legal"
    
    return "autre"


# ═══════════════════════════════════════════════════════════════
# MAIN ENDPOINT
# ═══════════════════════════════════════════════════════════════
@app.post("/parse")
async def parse_document(
    file: UploadFile = File(...),
    authorization: Optional[str] = Header(None),
):
    """Parse any document and return structured text."""
    
    # Simple auth check
    if authorization != f"Bearer {PARSER_API_KEY}":
        raise HTTPException(status_code=401, detail="Invalid API key")
    
    # Read file
    file_bytes = await file.read()
    if len(file_bytes) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail=f"File too large (max {MAX_FILE_SIZE // 1024 // 1024}MB)")
    
    if len(file_bytes) == 0:
        raise HTTPException(status_code=400, detail="Empty file")
    
    filename = file.filename or "unknown"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    
    print(f"[parser] Processing: {filename} ({len(file_bytes)} bytes, ext={ext})")
    
    try:
        result = None
        
        # Route to the right parser
        if ext == "pdf":
            result = parse_pdf(file_bytes, filename)
            # If pymupdf found almost no text, try OCR
            if result["quality"] in ("low", "failed"):
                print(f"[parser] PDF text extraction poor ({len(result['content'])} chars), trying OCR...")
                try:
                    ocr_result = parse_pdf_ocr(file_bytes, filename)
                    if len(ocr_result["content"]) > len(result["content"]):
                        result = ocr_result
                except Exception as e:
                    print(f"[parser] OCR failed: {e}")
        
        elif ext in ("xlsx", "xls"):
            result = parse_excel(file_bytes, filename)
        
        elif ext in ("docx", "doc"):
            result = parse_docx(file_bytes, filename)
        
        elif ext in ("pptx", "ppt"):
            result = parse_pptx(file_bytes, filename)
        
        elif ext in ("jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff"):
            result = parse_image(file_bytes, filename)
        
        elif ext in ("csv", "tsv", "txt", "md"):
            result = parse_text_file(file_bytes, filename, ext)
        
        else:
            # Try as text
            try:
                result = parse_text_file(file_bytes, filename, "txt")
            except:
                raise HTTPException(status_code=400, detail=f"Unsupported format: .{ext}")
        
        # Classify document type
        category = classify_document(filename, result.get("content", ""))
        
        # Build summary
        content_len = len(result.get("content", ""))
        parts = [result.get("method", "unknown")]
        if result.get("pages"):
            parts.append(f"{result['pages']} pages")
        if result.get("sheets"):
            parts.append(f"{result['sheets']} feuilles")
        if result.get("slides"):
            parts.append(f"{result['slides']} slides")
        if result.get("tables_found"):
            parts.append(f"{result['tables_found']} tableaux")
        parts.append(f"{content_len:,} caractères")
        summary = " — ".join(parts)
        
        print(f"[parser] Done: {filename} → {summary}")
        
        return JSONResponse({
            "success": True,
            "fileName": filename,
            "content": result.get("content", ""),
            "method": result.get("method", "unknown"),
            "category": category,
            "quality": result.get("quality", "unknown"),
            "summary": summary,
            "pages": result.get("pages"),
            "sheets": result.get("sheets"),
            "slides": result.get("slides"),
            "tables_found": result.get("tables_found", 0),
            "chars_extracted": content_len,
        })
    
    except HTTPException:
        raise
    except Exception as e:
        print(f"[parser] Error processing {filename}: {traceback.format_exc()}")
        return JSONResponse({
            "success": False,
            "fileName": filename,
            "content": f"[Erreur de parsing: {str(e)}]",
            "method": "error",
            "category": "autre",
            "quality": "failed",
            "summary": f"Erreur: {str(e)[:200]}",
            "chars_extracted": 0,
        }, status_code=200)  # Return 200 so the frontend can continue with other files


@app.get("/health")
async def health():
    """Health check endpoint."""
    return {"status": "ok", "version": "1.2.0", "features": ["parse", "ovo-excel", "odd-excel"]}


# ── Register Excel generation endpoints ──
from ovo_excel import register_excel_endpoints
register_excel_endpoints(app)

from odd_excel import register_odd_endpoints
register_odd_endpoints(app)


if __name__ == "__main__":
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
