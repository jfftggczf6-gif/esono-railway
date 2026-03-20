"""
ESONO Investment Memo PPTX Generator — Python micro-service endpoint
Creates a professional 20+ slide presentation from memo JSON data.

Palette: Navy #0F2B46, Blue #1B5E8A, Teal #0E7C6B, Gold #C4841D, Red #9B2C2C
Fonts: Georgia (titles), Calibri (body)
"""

import io
import base64
from typing import Optional
from fastapi import HTTPException, Header, Request
from fastapi.responses import Response
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ═══ Colors ═══
NAVY = RGBColor(0x0F, 0x2B, 0x46)
BLUE = RGBColor(0x1B, 0x5E, 0x8A)
TEAL = RGBColor(0x0E, 0x7C, 0x6B)
GOLD = RGBColor(0xC4, 0x84, 0x1D)
RED = RGBColor(0x9B, 0x2C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
BLACK = RGBColor(0x33, 0x33, 0x33)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.63)

safe = lambda v, fb="—": str(v) if v not in (None, "", []) else fb
arr = lambda v: v if isinstance(v, list) else []
trunc = lambda s, mx=280: (s[:mx] + "…" if s and len(s) > mx else s or "")


def add_footer(slide, prs, company, num):
    """Add footer bar to content slide."""
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.33), SLIDE_W, Inches(0.3)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = LIGHT
    slide.shapes[-1].line.fill.background()
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(5.33), Inches(8), Inches(0.3)).text_frame
    tf.paragraphs[0].text = f"ESONO — Investment Memorandum — {company} — Confidentiel"
    tf.paragraphs[0].font.size = Pt(7)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = "Calibri"
    tf2 = slide.shapes.add_textbox(Inches(9), Inches(5.33), Inches(0.8), Inches(0.3)).text_frame
    tf2.paragraphs[0].text = str(num)
    tf2.paragraphs[0].font.size = Pt(7)
    tf2.paragraphs[0].font.color.rgb = GRAY
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_title_bar(slide, title):
    """Add navy title bar at top of content slide."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = NAVY
    shp.line.fill.background()
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5)).text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.name = "Georgia"
    p.font.bold = True


def add_body_text(slide, text, left=0.5, top=1.0, width=9, height=4.0, size=11):
    """Add body text block."""
    tf = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = trunc(str(text), 1500)
    p.font.size = Pt(size)
    p.font.color.rgb = BLACK
    p.font.name = "Calibri"
    p.line_spacing = Pt(16)
    return tf


def add_bullets(slide, items, left=0.5, top=1.0, width=9, height=4.0, size=11):
    """Add bulleted list."""
    tf = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items[:12]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {trunc(str(item), 200)}"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return tf


def add_kpi_card(slide, x, y, value, label, w=2.1, h=1.0, color=NAVY):
    """Add a KPI card with big number and label."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = LIGHT
    shp.line.color.rgb = BORDER
    shp.line.width = Pt(1)
    # Value
    tf = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(w), Inches(0.5)).text_frame
    tf.paragraphs[0].text = safe(value)
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.name = "Georgia"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Label
    tf2 = slide.shapes.add_textbox(Inches(x), Inches(y + 0.6), Inches(w), Inches(0.3)).text_frame
    tf2.paragraphs[0].text = label
    tf2.paragraphs[0].font.size = Pt(8)
    tf2.paragraphs[0].font.color.rgb = GRAY
    tf2.paragraphs[0].font.name = "Calibri"
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_table(slide, headers, rows, left=0.5, top=1.0, width=9.0, row_h=0.35):
    """Add a formatted table."""
    n_rows = min(len(rows) + 1, 12)
    n_cols = len(headers)
    col_w = width / n_cols
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(row_h * n_rows)).table
    # Headers
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.font.bold = True
    # Data
    for i, row in enumerate(rows[:11]):
        for j, val in enumerate(row[:n_cols]):
            cell = tbl.cell(i + 1, j)
            cell.text = safe(val, "")
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = BLACK
                p.font.name = "Calibri"
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT


def generate_memo_pptx(data: dict) -> io.BytesIO:
    """Generate the Investment Memo PPTX from JSON data."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # Blank layout
    
    company = safe(data.get("page_de_garde", {}).get("titre", ""), "Investment Memorandum")
    company_short = safe(data.get("presentation_entreprise", {}).get("historique", "")[:50], "Entreprise")
    sn = [0]  # slide number counter
    
    def ns():
        sn[0] += 1
        return sn[0]

    # ═══ SLIDE 1: Cover ═══
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    # Decorative circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-0.5), Inches(4), Inches(4))
    circ.fill.background()
    circ.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    circ.line.width = Pt(2)
    circ.line.color.brightness = 0.9
    # Badge
    tf = s.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(3), Inches(0.3)).text_frame
    tf.paragraphs[0].text = "CONFIDENTIEL"
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"
    # Title
    pg = data.get("page_de_garde", {})
    tf = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7), Inches(1.2)).text_frame
    tf.paragraphs[0].text = "Investment Memorandum"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Georgia"
    tf.paragraphs[0].font.bold = True
    # Subtitle
    tf = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(7), Inches(0.8)).text_frame
    tf.paragraphs[0].text = safe(pg.get("titre", ""))
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"
    # Bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(1), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()
    # Meta
    meta = f"Date: {safe(pg.get('date', ''))}\nVersion: {safe(pg.get('version', 'v1.0'))}\nPréparé par: ESONO"
    tf = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(5), Inches(1)).text_frame
    tf.paragraphs[0].text = meta
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].line_spacing = Pt(18)

    # ═══ SLIDE 2: TOC ═══
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Table des matières")
    sections = ["Résumé exécutif", "Présentation entreprise", "Analyse de marché",
                "Modèle économique", "Analyse financière", "Projections financières",
                "Valorisation", "Besoins de financement", "Équipe & Gouvernance",
                "ESG & Impact", "Analyse des risques", "Thèse d'investissement",
                "Structure proposée", "Recommandation finale"]
    for i, sec in enumerate(sections):
        y = 1.0 + i * 0.28
        tf = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(0.5), Inches(0.25)).text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = BLUE
        tf.paragraphs[0].font.bold = True
        tf2 = s.shapes.add_textbox(Inches(1.3), Inches(y), Inches(7), Inches(0.25)).text_frame
        tf2.paragraphs[0].text = sec
        tf2.paragraphs[0].font.size = Pt(10)
        tf2.paragraphs[0].font.color.rgb = BLACK
        tf2.paragraphs[0].font.name = "Calibri"
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 3: Résumé exécutif ═══
    re = data.get("resume_executif", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "1. Résumé exécutif")
    # Score + Verdict
    verdict = safe(re.get("recommandation_preliminaire", ""))
    score = safe(re.get("score_ir", ""))
    vcolor = TEAL if "INVESTIR" in verdict.upper() else (GOLD if "APPROFONDIR" in verdict.upper() else RED)
    add_kpi_card(s, 0.5, 1.0, f"{score}/100", "Score IR", 2.0, 0.9, NAVY)
    add_kpi_card(s, 2.7, 1.0, verdict, "Recommandation", 2.5, 0.9, vcolor)
    # Points clés
    add_bullets(s, arr(re.get("points_cles", [])), 0.5, 2.2, 9, 2.5, 10)
    add_footer(s, prs, company_short, ns())

    # Résumé narratif
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "1. Résumé exécutif — Synthèse")
    add_body_text(s, safe(re.get("synthese", "")), size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 5: Présentation entreprise ═══
    pe = data.get("presentation_entreprise", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "2. Présentation de l'entreprise")
    add_body_text(s, safe(pe.get("historique", "")), height=2.0, size=10)
    add_body_text(s, f"Activités: {safe(pe.get('activites', ''))}\n\nGouvernance: {safe(pe.get('gouvernance', ''))}", top=3.2, height=2.0, size=9)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 6: Marché ═══
    am = data.get("analyse_marche", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "3. Analyse de marché")
    add_body_text(s, f"Taille du marché: {safe(am.get('taille_marche', ''))}\n\nDynamiques: {safe(am.get('dynamiques', ''))}\n\nConcurrence: {safe(am.get('concurrence', ''))}", size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 7: Modèle économique ═══
    me = data.get("modele_economique", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "4. Modèle économique")
    add_body_text(s, f"Sources de revenus: {safe(me.get('sources_revenus', ''))}\n\nStructure de coûts: {safe(me.get('structure_couts', ''))}\n\nAvantages concurrentiels: {safe(me.get('avantages_concurrentiels', ''))}", size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 8-9: Analyse financière ═══
    af = data.get("analyse_financiere", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "5. Analyse financière — KPIs")
    kpis = arr(af.get("ratios_cles", []))
    for i, kpi in enumerate(kpis[:8]):
        x = 0.5 + (i % 4) * 2.3
        y = 1.0 + (i // 4) * 1.3
        label = kpi.get("ratio", kpi.get("nom", "")) if isinstance(kpi, dict) else str(kpi)
        value = kpi.get("valeur", "") if isinstance(kpi, dict) else ""
        add_kpi_card(s, x, y, safe(value), safe(label), 2.1, 1.0)
    add_footer(s, prs, company_short, ns())

    s = prs.slides.add_slide(blank)
    add_title_bar(s, "5. Analyse financière — Commentaire")
    add_body_text(s, safe(af.get("commentaire_global", af.get("analyse_narrative", ""))), size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 10: Projections ═══
    proj = data.get("projections_financieres", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "6. Projections financières 5 ans")
    tableau = arr(proj.get("tableau", proj.get("projections", [])))
    if tableau:
        headers = list(tableau[0].keys()) if isinstance(tableau[0], dict) else ["Poste", "An1", "An2", "An3", "An4", "An5"]
        rows = [list(row.values()) if isinstance(row, dict) else row for row in tableau[:10]]
        add_table(s, headers, rows, top=1.0)
    else:
        add_body_text(s, safe(proj.get("synthese", "Projections non disponibles")))
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 11: Valorisation ═══
    val = data.get("valorisation", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "7. Valorisation")
    add_kpi_card(s, 0.5, 1.0, safe(val.get("fourchette_valorisation", val.get("valeur_estimee", ""))), "Fourchette de valorisation", 4.0, 1.0, TEAL)
    methods = arr(val.get("methodes", []))
    if methods:
        for i, m in enumerate(methods[:3]):
            y = 2.3 + i * 0.8
            name = m.get("methode", m.get("nom", "")) if isinstance(m, dict) else str(m)
            value = m.get("valeur", "") if isinstance(m, dict) else ""
            tf = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3), Inches(0.3)).text_frame
            tf.paragraphs[0].text = f"• {safe(name)}: {safe(value)}"
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.name = "Calibri"
    add_body_text(s, safe(val.get("commentaire", "")), top=4.2, height=0.8, size=9)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 12: Besoins de financement ═══
    bf = data.get("besoins_financement", data.get("structure_proposee", {}))
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "8. Besoins de financement")
    add_kpi_card(s, 0.5, 1.0, safe(bf.get("montant", bf.get("montant_total", ""))), "Montant recherché", 3.0, 1.0, GOLD)
    add_body_text(s, f"Instrument: {safe(bf.get('instrument', ''))}\n\nUtilisation: {safe(bf.get('utilisation', bf.get('utilisation_fonds', '')))}", top=2.3, size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 13: Équipe ═══
    eq = data.get("equipe_gouvernance", data.get("presentation_entreprise", {}))
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "9. Équipe & Gouvernance")
    add_body_text(s, f"Gouvernance: {safe(eq.get('gouvernance', ''))}\n\nEffectifs: {safe(eq.get('effectifs', ''))}", size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 14: ESG ═══
    esg = data.get("esg_impact", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "10. ESG & Impact")
    add_body_text(s, f"Impact social: {safe(esg.get('impact_social', ''))}\n\nImpact environnemental: {safe(esg.get('impact_environnemental', ''))}\n\nConformité IFC PS: {safe(esg.get('conformite_ifc_ps', ''))}", size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 15-16: Risques ═══
    ar = data.get("analyse_risques", {})
    risques = arr(ar.get("risques_identifies", []))
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "11. Matrice des risques")
    if risques:
        headers = ["Risque", "Probabilité", "Impact", "Mitigation"]
        rows = [[safe(r.get("description", r.get("categorie", "")))[:50], safe(r.get("probabilite", "")), safe(r.get("impact", "")), safe(r.get("mitigation", ""))[:60]] for r in risques[:8]]
        add_table(s, headers, rows)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 17: Thèse d'investissement ═══
    ti = data.get("these_investissement", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "12. Thèse d'investissement")
    # Two columns: positive vs negative
    add_body_text(s, f"POUR:\n{safe(ti.get('these_positive', ''))}", left=0.5, top=1.0, width=4.3, height=3.5, size=9)
    add_body_text(s, f"CONTRE:\n{safe(ti.get('these_negative', ''))}", left=5.2, top=1.0, width=4.3, height=3.5, size=9)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 18: Structure proposée ═══
    sp = data.get("structure_proposee", {})
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "13. Structure proposée")
    add_body_text(s, f"Instrument: {safe(sp.get('instrument', ''))}\nMontant: {safe(sp.get('montant', ''))}\nDilution: {safe(sp.get('dilution_estimee', ''))}", top=1.0, height=1.5, size=11)
    add_bullets(s, arr(sp.get("conditions_precedentes", [])), top=2.8, height=2.0, size=10)
    add_footer(s, prs, company_short, ns())

    # ═══ SLIDE 19: Recommandation finale ═══
    rf = data.get("recommandation_finale", {})
    s = prs.slides.add_slide(blank)
    # Dark background for final slide
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    verdict = safe(rf.get("verdict", ""))
    vcolor = TEAL if "INVESTIR" in verdict.upper() else (GOLD if "APPROFONDIR" in verdict.upper() else RED)
    tf = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.5)).text_frame
    tf.paragraphs[0].text = "RECOMMANDATION FINALE"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"
    # Big verdict
    tf = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.8)).text_frame
    tf.paragraphs[0].text = verdict
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = vcolor
    tf.paragraphs[0].font.name = "Georgia"
    tf.paragraphs[0].font.bold = True
    # Justification
    tf = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2.5)).text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = trunc(safe(rf.get("justification", "")), 800)
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].line_spacing = Pt(16)

    # ═══ SLIDE 20: Disclaimer ═══
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Disclaimer")
    disclaimer = """Ce document est strictement confidentiel et destiné uniquement aux personnes autorisées.

Les informations contenues dans ce mémorandum ont été préparées par ESONO sur la base des données fournies par l'entreprise et de sources considérées comme fiables. ESONO ne garantit pas l'exactitude ou l'exhaustivité de ces informations.

Ce document ne constitue pas une offre ou une sollicitation d'investissement. Toute décision d'investissement doit être prise après une due diligence approfondie.

Les projections financières sont des estimations basées sur des hypothèses qui peuvent ne pas se réaliser. Les performances passées ne préjugent pas des performances futures."""
    add_body_text(s, disclaimer, size=9)
    add_footer(s, prs, company_short, ns())

    # Save
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def register_pptx_endpoints(app):
    """Register PPTX generation endpoint."""

    @app.post("/generate-memo-pptx")
    async def gen_memo_pptx(request: Request, authorization: Optional[str] = Header(None)):
        import os
        PK = os.getenv("PARSER_API_KEY", "esono-parser-dev-key")
        if authorization != f"Bearer {PK}":
            raise HTTPException(status_code=401, detail="Invalid API key")

        body = await request.json()
        data = body.get("data", {})

        print(f"[memo-pptx] Generating for: {safe(data.get('page_de_garde', {}).get('titre', ''))}")

        try:
            output = generate_memo_pptx(data)
        except Exception as e:
            import traceback
            print(f"[memo-pptx] Error: {traceback.format_exc()}")
            raise HTTPException(status_code=500, detail=f"PPTX error: {e}")

        print(f"[memo-pptx] Done: {output.getbuffer().nbytes} bytes")

        return Response(
            content=output.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=Investment_Memo.pptx"}
        )
